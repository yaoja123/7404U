from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("/Users/jayden/Desktop/7404 comp/project/MaxSup_A0_HiFi_Poster.pptx")


PALETTE = {
    "bg": "0A162B",
    "bg2": "12284A",
    "card": "F5F1E8",
    "card_soft": "FBF8F1",
    "white": "FFFFFF",
    "ink": "10233F",
    "ink2": "28476E",
    "blue": "2C67F2",
    "blue_soft": "DDE8FF",
    "orange": "FF7A45",
    "orange_soft": "FFE3D7",
    "teal": "41B8A8",
    "teal_soft": "D9F3EE",
    "gold": "F6C453",
    "line": "C9D7EA",
    "muted": "6D84A3",
}


def rgb(hex_value: str) -> RGBColor:
    value = PALETTE.get(hex_value, hex_value).replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False, transparency=0.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1.4)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, x, y, w, h, fill="card", line="line", shadow=True):
    if shadow:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.13, y + 0.13, w, h, "06101F", None, transparency=0.55)
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    font="Aptos",
    size=20,
    color="ink",
    bold=False,
    align="left",
    margin=0.12,
    valign="top",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_label(slide, x, y, w, h, text, fill, text_color="white"):
    pill = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, None)
    pill.adjustments[0] = 0.35
    add_textbox(slide, x, y + 0.01, w, h - 0.02, text, size=14, color=text_color, bold=True, align="center", valign="middle")
    return pill


def add_rule(slide, x, y, w, h=0.05, fill="line", transparency=0.0):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, None, transparency=transparency)


def add_metric_tile(slide, x, y, w, h, title, main, sub, accent_fill, accent_text="white"):
    add_card(slide, x, y, w, h, fill="card_soft")
    add_label(slide, x + 0.25, y + 0.22, 2.25, 0.45, title, accent_fill, accent_text)
    add_textbox(slide, x + 0.2, y + 0.95, w - 0.4, 1.0, main, font="Aptos Display", size=33, color="ink", bold=True)
    add_textbox(slide, x + 0.2, y + 1.75, w - 0.4, h - 1.95, sub, size=15.5, color="muted")


def add_logit_bars(slide, x, y, w, h, accent, mode_text, penalize_idx):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "white", "line")
    add_textbox(slide, x + 0.15, y + 0.12, w - 0.3, 0.35, mode_text, size=14, color="ink2", bold=True, align="center")

    labels = ["cat", "fox", "dog"]
    values = [1.0, 2.8, 0.4]
    max_val = 3.0
    left = x + 0.4
    base_y = y + h - 0.55
    bar_w = 0.62
    gap = 0.45
    for idx, (label, value) in enumerate(zip(labels, values)):
        bar_h = 1.9 * (value / max_val)
        bx = left + idx * (bar_w + gap)
        add_shape(slide, MSO_SHAPE.RECTANGLE, bx, base_y - bar_h, bar_w, bar_h, accent if idx == penalize_idx else "blue_soft", None)
        add_textbox(slide, bx - 0.05, base_y + 0.05, bar_w + 0.1, 0.3, label, size=11.5, color="muted", align="center")
        add_textbox(slide, bx - 0.05, base_y - bar_h - 0.28, bar_w + 0.1, 0.25, f"{value:.1f}", size=11.5, color="ink2", bold=idx == penalize_idx, align="center")

    arrow_x = left + penalize_idx * (bar_w + gap) + 0.16
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, arrow_x, y + 0.55, 0.32, 0.55, accent, None)
    add_textbox(slide, x + 0.18, y + h - 0.25, w - 0.36, 0.3, "penalized target", size=11.5, color="muted", align="center")


def add_cluster_diagram(slide, x, y, collapsed=False):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, 2.7, 1.9, "teal_soft", "teal")
    add_shape(slide, MSO_SHAPE.OVAL, x + 2.4, y + 1.05, 2.3, 1.6, "orange_soft", "orange")
    points_a = [
        (0.65, 0.55),
        (0.92, 0.72),
        (1.22, 0.58),
        (1.08, 1.02),
        (0.7, 0.95),
        (1.38, 0.82),
    ]
    points_b = [
        (2.95, 1.55),
        (3.2, 1.82),
        (3.55, 1.6),
        (3.75, 2.1),
        (3.12, 2.18),
        (3.45, 1.35),
    ]
    if collapsed:
        points_a = [(0.98, 0.75), (1.05, 0.88), (1.12, 0.82), (0.92, 0.84), (1.0, 0.95), (1.16, 0.72)]
        points_b = [(3.3, 1.7), (3.38, 1.82), (3.46, 1.75), (3.25, 1.9), (3.34, 1.6), (3.42, 1.96)]
    for px, py in points_a:
        add_shape(slide, MSO_SHAPE.OVAL, x + px, y + py, 0.12, 0.12, "teal", None)
    for px, py in points_b:
        add_shape(slide, MSO_SHAPE.OVAL, x + px, y + py, 0.12, 0.12, "orange", None)


def build_poster() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(46.81)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 33.11, 46.81, "bg", None)
    add_shape(slide, MSO_SHAPE.OVAL, -4.0, -2.8, 12.0, 8.0, "blue", None, transparency=0.78)
    add_shape(slide, MSO_SHAPE.OVAL, 24.8, -2.0, 10.0, 7.2, "orange", None, transparency=0.88)
    add_shape(slide, MSO_SHAPE.OVAL, 18.8, 34.0, 13.0, 10.0, "blue", None, transparency=0.88)

    for idx in range(11):
        add_rule(slide, 0.9, 5.0 + idx * 3.6, 31.1, 0.025, fill="white", transparency=0.92)
    for idx in range(7):
        add_rule(slide, 4.2 + idx * 4.1, 0.8, 0.025, 45.0, fill="white", transparency=0.95)

    add_label(slide, 1.05, 0.85, 4.45, 0.55, "A0 EDITABLE POWERPOINT", "blue")
    add_textbox(slide, 1.0, 1.65, 16.8, 1.6, "MaxSup", font="Aptos Display", size=68, color="white", bold=True)
    add_textbox(
        slide,
        1.0,
        3.05,
        16.6,
        2.3,
        "Fixing Label Smoothing When the Model Is Wrong",
        font="Aptos Display",
        size=30,
        color="white",
        bold=True,
    )
    add_textbox(
        slide,
        1.05,
        5.35,
        15.8,
        2.4,
        "A conference-booth poster concept that explains one precise idea:\nLabel Smoothing penalizes the true class; MaxSup suppresses the dominant prediction.",
        size=18.5,
        color="blue_soft",
    )
    add_textbox(slide, 1.1, 8.25, 9.5, 0.9, "WHEN PREDICTION IS WRONG", size=14.5, color="orange", bold=True)
    add_textbox(
        slide,
        1.0,
        8.9,
        12.5,
        2.25,
        "LS keeps pushing down the ground-truth logit.\nMaxSup directly suppresses the top-1 logit instead.",
        font="Aptos Display",
        size=25,
        color="white",
        bold=True,
    )

    hero = add_card(slide, 18.8, 1.0, 13.0, 10.9, fill="card")
    add_label(slide, 19.25, 1.45, 3.2, 0.5, "HERO COMPARISON", "ink")
    add_textbox(slide, 19.3, 2.15, 11.8, 1.1, "Wrong prediction case", font="Aptos Display", size=25, color="ink", bold=True)
    add_textbox(slide, 19.35, 2.95, 5.4, 1.3, "GT = cat\ncat = 1.0\nfox = 2.8\ndog = 0.4", size=18, color="ink2")
    add_textbox(slide, 24.3, 2.95, 6.6, 1.1, "A single substitution changes the training signal:\nz_gt  ->  z_max", size=17.5, color="ink2")
    add_logit_bars(slide, 19.4, 5.0, 5.6, 4.4, "orange", "LS penalizes z_gt", 0)
    add_logit_bars(slide, 25.1, 5.0, 5.6, 4.4, "blue", "MaxSup penalizes z_max", 1)
    add_textbox(slide, 19.55, 9.7, 11.0, 1.0, "Takeaway: MaxSup preserves LS regularization while avoiding error amplification.", size=16.5, color="muted", bold=True)

    add_card(slide, 1.0, 12.4, 31.05, 1.35, fill="bg2", line=None, shadow=False)
    add_textbox(
        slide,
        1.45,
        12.65,
        30.0,
        0.7,
        "LS regularizes the label target.   MaxSup regularizes the dominant prediction.",
        font="Aptos Display",
        size=23,
        color="white",
        bold=True,
        align="center",
        valign="middle",
    )

    add_textbox(slide, 18.65, 14.35, 12.7, 1.1, "z_gt  ->  z_max", font="Aptos Display", size=44, color="blue_soft", bold=True, align="right")

    left_panel = add_card(slide, 1.0, 14.2, 18.6, 13.15, fill="card")
    add_label(slide, 1.4, 14.65, 2.6, 0.48, "PROBLEM", "orange")
    add_textbox(
        slide,
        1.42,
        15.35,
        17.45,
        2.1,
        "Label Smoothing helps calibration, but on misclassified samples it can amplify the wrong direction and compress features too aggressively.",
        font="Aptos Display",
        size=24,
        color="ink",
        bold=True,
    )
    add_rule(slide, 1.45, 17.55, 17.7, 0.04, fill="line")
    add_label(slide, 1.4, 18.05, 2.45, 0.48, "METHOD", "blue")
    add_textbox(slide, 1.42, 18.75, 8.8, 1.6, "L_LS = alpha (z_gt - mean(z))", font="Georgia", size=24, color="ink", bold=True)
    add_textbox(slide, 1.42, 20.1, 8.8, 1.6, "L_MaxSup = alpha (z_max - mean(z))", font="Georgia", size=24, color="ink", bold=True)
    add_textbox(slide, 9.85, 18.8, 8.4, 2.6, "One change only:\nz_gt -> z_max", font="Aptos Display", size=24, color="blue", bold=True, align="center")
    add_card(slide, 1.45, 22.2, 17.3, 4.35, fill="card_soft", line="line", shadow=False)
    add_textbox(
        slide,
        1.75,
        22.55,
        16.7,
        3.7,
        "Why this matters\n\n• Correct prediction: still suppresses overconfidence.\n• Wrong prediction: suppresses the dominant wrong class instead of hurting the truth.\n• Representation space: preserves richer intra-class variation while keeping classes separable.",
        size=17.5,
        color="ink2",
    )

    right_panel = add_card(slide, 20.3, 14.2, 11.75, 13.15, fill="bg2", line="blue")
    add_textbox(slide, 20.85, 14.75, 10.6, 0.85, "RESULTS THAT STOP PEOPLE", font="Aptos Display", size=25, color="white", bold=True)
    add_metric_tile(slide, 20.8, 15.95, 5.0, 3.9, "DeiT-Small", "76.49", "MaxSup beats LS (76.08) and OLS (76.16).\nMessage: not CNN-specific; works for transformers too.", "blue")
    add_metric_tile(slide, 26.1, 15.95, 5.2, 3.9, "ADE20K", "42.8", "UPerNet segmentation also improves over LS (42.4).\nMessage: better backbone transfer, not just classification.", "teal", accent_text="ink")
    add_metric_tile(slide, 20.8, 20.25, 5.0, 3.4, "Grad-CAM", "Sharper focus", "More object-centric attention.\nLess background and shortcut bias.", "orange", accent_text="ink")
    add_metric_tile(slide, 26.1, 20.25, 5.2, 3.4, "Overhead", "~ none", "Appendix Table 16: negligible compute cost.\nSimple swap, no architecture change.", "gold", accent_text="ink")
    add_card(slide, 20.8, 24.05, 10.5, 2.7, fill="card")
    add_label(slide, 21.1, 24.35, 2.0, 0.43, "DEMO", "ink")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 27.9, 24.45, 2.35, 1.85, "card_soft", "muted")
    add_textbox(slide, 21.15, 24.95, 5.9, 1.25, "Interactive logit simulator\nScan to test LS vs MaxSup.", size=18, color="ink", bold=True)
    add_textbox(slide, 28.15, 25.0, 1.85, 1.1, "QR", font="Aptos Display", size=28, color="muted", bold=True, align="center", valign="middle")

    add_textbox(slide, 1.0, 28.25, 18.0, 0.9, "SUPPORT EVIDENCE", font="Aptos Display", size=28, color="white", bold=True)

    card_y = 29.15
    card_h = 12.7
    add_card(slide, 1.0, card_y, 10.1, card_h, fill="card")
    add_label(slide, 1.35, card_y + 0.35, 3.15, 0.46, "REPRESENTATION", "teal", text_color="ink")
    add_textbox(slide, 1.35, card_y + 1.05, 9.2, 1.2, "MaxSup improves inter-class separability and preserves intra-class variability.", font="Aptos Display", size=21, color="ink", bold=True)
    add_cluster_diagram(slide, 1.65, card_y + 2.8, collapsed=True)
    add_textbox(slide, 1.55, card_y + 5.7, 4.0, 0.45, "LS: clusters collapse", size=14.5, color="orange", bold=True, align="center")
    add_cluster_diagram(slide, 5.25, card_y + 2.8, collapsed=False)
    add_textbox(slide, 5.15, card_y + 5.7, 4.1, 0.45, "MaxSup: richer spread", size=14.5, color="teal", bold=True, align="center")
    add_textbox(slide, 1.45, card_y + 6.4, 8.9, 4.6, "Poster note\n\nUse Table 2 to say: same class should not collapse too much, different classes should stay clearly separated.\n\nThis gives you an easy visual explanation for why the method helps downstream tasks.", size=16.2, color="ink2")

    add_card(slide, 11.55, card_y, 10.1, card_h, fill="card")
    add_label(slide, 11.9, card_y + 0.35, 3.5, 0.46, "DOWNSTREAM", "blue")
    add_textbox(slide, 11.9, card_y + 1.05, 9.2, 1.2, "MaxSup is not just an ImageNet trick. It supports transfer, fine-grained recognition, long-tailed data, and segmentation.", font="Aptos Display", size=20.5, color="ink", bold=True)
    add_metric_tile(slide, 12.0, card_y + 2.65, 4.4, 2.65, "Fine-grained", "CUB / Cars", "Captures subtle details better.", "orange", accent_text="ink")
    add_metric_tile(slide, 16.55, card_y + 2.65, 4.4, 2.65, "Long-tailed", "better tradeoff", "More robust across many-, medium-, and low-shot regimes.", "teal", accent_text="ink")
    add_metric_tile(slide, 12.0, card_y + 5.65, 4.4, 2.65, "OOD / Corruption", "reliable confidence", "Better calibration on CIFAR-10-C.", "blue")
    add_metric_tile(slide, 16.55, card_y + 5.65, 4.4, 2.65, "Transfer", "stronger backbone", "Segmentation results show features transfer better.", "gold", accent_text="ink")
    add_textbox(slide, 11.95, card_y + 8.8, 8.9, 2.5, "Booth line\n\nThis small loss change improves accuracy, transferability, segmentation, calibration, and interpretability without changing the model architecture.", size=16.2, color="ink2")

    add_card(slide, 22.1, card_y, 9.95, card_h, fill="card")
    add_label(slide, 22.45, card_y + 0.35, 3.0, 0.46, "APPENDIX", "ink")
    add_textbox(slide, 22.45, card_y + 1.05, 8.95, 1.2, "Use appendix results strategically rather than dumping them.", font="Aptos Display", size=21, color="ink", bold=True)
    add_metric_tile(slide, 22.55, card_y + 2.65, 8.95, 2.15, "Table 12", "Extended transfer learning", "Very poster-friendly. Strong support for representation quality.", "blue")
    add_metric_tile(slide, 22.55, card_y + 5.1, 8.95, 2.15, "Table 13", "More baselines", "Useful when someone asks whether the comparison is broad enough.", "orange", accent_text="ink")
    add_metric_tile(slide, 22.55, card_y + 7.55, 8.95, 2.15, "Table 16", "Compute efficiency", "The method is lightweight. Almost no extra cost.", "teal", accent_text="ink")
    add_textbox(slide, 22.6, card_y + 10.1, 8.8, 1.4, "Rule of thumb\nPick only 1-2 appendix items in the final version. Curate; do not crowd.", size=16.2, color="ink2", bold=True)

    add_card(slide, 1.0, 42.45, 31.05, 2.7, fill="card")
    add_textbox(
        slide,
        1.35,
        42.82,
        18.8,
        0.75,
        "simple  ·  no architecture change  ·  low overhead  ·  stronger representations",
        font="Aptos Display",
        size=22,
        color="ink",
        bold=True,
    )
    add_textbox(
        slide,
        1.35,
        43.55,
        19.2,
        0.9,
        "Editable PPT handoff: all major blocks are shapes + text boxes so your team can restyle, resize, and replace content directly in PowerPoint.",
        size=15.5,
        color="muted",
    )
    add_textbox(slide, 22.0, 42.85, 9.0, 0.65, "Speaking split", size=15.5, color="orange", bold=True, align="right")
    add_textbox(slide, 22.0, 43.4, 9.0, 0.9, "1 overview   2 theory   3 results   4 demo + Q&A", size=15.5, color="ink2", align="right")

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    result = build_poster()
    print(result)
