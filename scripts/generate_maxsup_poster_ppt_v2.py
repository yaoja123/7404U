from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("/Users/jayden/Desktop/7404 comp/project/MaxSup_A0_HiFi_Poster_v2.pptx")

PALETTE = {
    "bg": "FBF8F1",
    "paper": "FFFDF9",
    "sand": "F4E8D8",
    "sand2": "F8EFE4",
    "blue_soft": "EAF1FB",
    "blue_mid": "CFE0F7",
    "blue": "2458B3",
    "ink": "1D2A39",
    "ink2": "42576E",
    "muted": "6E7E8F",
    "orange": "E76F51",
    "orange_soft": "FDE6DF",
    "teal": "2A7F62",
    "teal_soft": "DDF2EA",
    "gold": "E6B655",
    "line": "D6DEE7",
    "white": "FFFFFF",
}


def rgb(name_or_hex: str) -> RGBColor:
    value = PALETTE.get(name_or_hex, name_or_hex).replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_shape(slide, kind, x, y, w, h, fill, line=None, transparency=0.0):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, x, y, w, h, fill="paper", line="line", shadow=True):
    if shadow:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.11, y + 0.12, w, h, "DCE5EF", None, transparency=0.45)
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=20,
    color="ink",
    bold=False,
    font="Aptos",
    align="left",
    valign="top",
    margin=0.10,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    r = p.runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def add_pill(slide, x, y, w, h, text, fill, text_color="white"):
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, None)
    shp.adjustments[0] = 0.38
    add_textbox(slide, x, y + 0.01, w, h - 0.02, text, size=14, color=text_color, bold=True, align="center", valign="middle")
    return shp


def add_rule(slide, x, y, w, h=0.03, fill="line"):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, None)


def add_bar_group(slide, x, y, title, accent, target_label, target_idx):
    add_textbox(slide, x, y, 3.9, 0.45, title, size=20, color="ink", bold=True, align="center")
    labels = ["cat", "fox", "dog"]
    values = [1.0, 2.8, 0.4]
    max_val = 3.0
    base_y = y + 3.0
    bar_w = 0.62
    gap = 0.42
    start_x = x + 0.9
    for i, (lab, val) in enumerate(zip(labels, values)):
        height = 1.95 * (val / max_val)
        bx = start_x + i * (bar_w + gap)
        add_shape(slide, MSO_SHAPE.RECTANGLE, bx, base_y - height, bar_w, height, accent if i == target_idx else "blue_mid", None)
        add_textbox(slide, bx - 0.05, base_y + 0.05, bar_w + 0.1, 0.25, lab, size=12.5, color="muted", align="center")
        add_textbox(slide, bx - 0.1, base_y - height - 0.28, bar_w + 0.2, 0.22, f"{val:.1f}", size=12.5, color="ink2", bold=i == target_idx, align="center")
    arrow = add_shape(slide, MSO_SHAPE.DOWN_ARROW, start_x + target_idx * (bar_w + gap) + 0.15, y + 0.65, 0.32, 0.58, accent, None)
    add_textbox(slide, x + 0.1, y + 3.35, 3.7, 0.28, target_label, size=13, color="ink2", align="center")
    return arrow


def add_result_tile(slide, x, y, w, h, label, main, note, fill, label_color="white"):
    add_card(slide, x, y, w, h, fill="paper")
    add_pill(slide, x + 0.24, y + 0.22, 2.0, 0.42, label, fill, label_color)
    add_textbox(slide, x + 0.2, y + 0.78, w - 0.4, 0.72, main, size=30, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, x + 0.2, y + 1.55, w - 0.4, h - 1.75, note, size=17, color="ink2")


def add_small_card(slide, x, y, w, h, title, body, fill):
    add_card(slide, x, y, w, h, fill=fill, shadow=False)
    add_textbox(slide, x + 0.18, y + 0.16, w - 0.36, 0.5, title, size=18, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, x + 0.18, y + 0.72, w - 0.36, h - 0.9, body, size=15.5, color="ink2")


def add_cluster_icon(slide, x, y, color_fill, color_dot, spread=True):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, 2.8, 1.75, color_fill, None)
    pts = [(0.55, 0.45), (0.92, 0.7), (1.22, 0.56), (1.38, 0.9), (0.82, 1.05), (1.58, 0.62)]
    if not spread:
        pts = [(1.0, 0.76), (1.08, 0.84), (0.95, 0.9), (1.14, 0.72), (1.03, 0.98), (1.16, 0.86)]
    for px, py in pts:
        add_shape(slide, MSO_SHAPE.OVAL, x + px, y + py, 0.11, 0.11, color_dot, None)


def build():
    prs = Presentation()
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(46.81)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 33.11, 46.81, "bg", None)
    add_shape(slide, MSO_SHAPE.OVAL, 22.8, -1.0, 8.0, 4.2, "orange_soft", None)
    add_shape(slide, MSO_SHAPE.OVAL, -1.8, 9.8, 6.5, 4.0, "blue_soft", None)
    add_shape(slide, MSO_SHAPE.OVAL, 24.5, 31.0, 6.5, 5.0, "teal_soft", None)

    add_pill(slide, 1.0, 0.95, 2.5, 0.48, "A0 PORTRAIT", "blue")
    add_textbox(slide, 1.0, 1.75, 20.0, 1.1, "MaxSup", size=62, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 1.0, 2.88, 20.8, 1.65, "Fixing Label Smoothing\nWhen the Model Is Wrong", size=33, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 1.0, 4.68, 17.5, 0.95, "One idea. One substitution. Better behavior when the model is wrong.", size=21, color="ink2")

    add_card(slide, 21.55, 1.0, 10.45, 3.8, fill="paper")
    add_pill(slide, 21.9, 1.35, 2.45, 0.42, "GROUP INFO", "sand", text_color="ink")
    add_textbox(slide, 21.9, 1.92, 9.7, 0.72, "Group ID: ________", size=19, color="ink", bold=True)
    add_textbox(slide, 21.9, 2.55, 9.7, 1.1, "Names: ________ / ________ / ________ / ________", size=18, color="ink2")
    add_textbox(slide, 21.9, 3.45, 9.7, 0.5, "Put your real names here for the final submission.", size=14.5, color="muted")

    hero = add_card(slide, 1.0, 6.0, 31.1, 11.0, fill="paper")
    add_pill(slide, 1.45, 6.4, 3.25, 0.48, "EYE-CATCHING MAIN GRAPHIC", "orange", "white")
    add_textbox(slide, 1.45, 7.1, 12.0, 0.95, "Wrong prediction example", size=28, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 1.45, 8.0, 6.2, 2.2, "GT = cat\ncat = 1.0\nfox = 2.8\ndog = 0.4", size=25, color="ink2", bold=True)
    add_textbox(slide, 1.45, 10.55, 8.2, 1.25, "LS penalizes the true class\nwhen the model is already wrong.", size=20.5, color="orange", bold=True)
    add_textbox(slide, 1.45, 12.2, 8.3, 1.25, "MaxSup suppresses the dominant\nwrong logit instead.", size=20.5, color="blue", bold=True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.15, 7.3, 10.25, 7.95, "sand2", None)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 20.7, 7.3, 10.0, 7.95, "blue_soft", None)
    add_bar_group(slide, 12.65, 8.1, "LS", "orange", "penalize z_gt", 0)
    add_bar_group(slide, 23.15, 8.1, "MaxSup", "blue", "penalize z_max", 1)
    add_textbox(slide, 18.6, 10.0, 3.3, 1.2, "z_gt\n->\nz_max", size=26, color="ink", bold=True, font="Aptos Display", align="center", valign="middle")
    add_shape(slide, MSO_SHAPE.CHEVRON, 18.85, 8.15, 2.8, 1.0, "white", None)
    add_shape(slide, MSO_SHAPE.CHEVRON, 18.85, 13.0, 2.8, 1.0, "white", None)
    add_textbox(slide, 10.55, 14.6, 9.25, 0.65, "LS regularizes the label target", size=18.5, color="ink2", bold=True, align="center")
    add_textbox(slide, 21.15, 14.6, 9.0, 0.65, "MaxSup regularizes the dominant prediction", size=18.5, color="ink2", bold=True, align="center")

    add_textbox(slide, 1.0, 17.7, 10.0, 0.7, "WHY THIS MATTERS", size=29, color="ink", bold=True, font="Aptos Display")
    add_rule(slide, 1.0, 18.55, 31.0, 0.03)

    add_card(slide, 1.0, 19.1, 9.85, 6.0, fill="paper")
    add_pill(slide, 1.3, 19.4, 1.85, 0.42, "PROBLEM", "orange")
    add_textbox(slide, 1.3, 20.05, 9.1, 1.45, "Label Smoothing helps calibration, but can amplify errors on misclassified samples.", size=23, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 1.3, 21.8, 9.0, 2.4, "People will not read a paragraph here.\nKeep this block to the single failure mode you want them to remember.", size=18, color="ink2")

    add_card(slide, 11.6, 19.1, 9.85, 6.0, fill="paper")
    add_pill(slide, 11.9, 19.4, 1.8, 0.42, "METHOD", "blue")
    add_textbox(slide, 11.9, 20.0, 8.9, 1.0, "One substitution", size=24, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 11.9, 20.95, 8.9, 0.8, "z_gt  ->  z_max", size=30, color="blue", bold=True, font="Aptos Display")
    add_textbox(slide, 11.9, 22.0, 8.9, 1.5, "L_LS = alpha (z_gt - mean(z))\nL_MaxSup = alpha (z_max - mean(z))", size=20, color="ink2")
    add_textbox(slide, 11.9, 23.65, 8.9, 0.7, "Simple change. Consistent signal.", size=18, color="muted", bold=True)

    add_card(slide, 22.2, 19.1, 9.9, 6.0, fill="paper")
    add_pill(slide, 22.5, 19.4, 2.15, 0.42, "TAKEAWAY", "teal", "white")
    add_textbox(slide, 22.5, 20.05, 9.0, 1.5, "Suppress the overconfident wrong class, not the ground truth.", size=23, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 22.5, 21.95, 8.95, 2.15, "This keeps regularization on correct cases and avoids error amplification on wrong ones.", size=18, color="ink2")

    add_textbox(slide, 1.0, 26.0, 10.0, 0.7, "MAIN RESULTS", size=29, color="ink", bold=True, font="Aptos Display")
    add_rule(slide, 1.0, 26.82, 31.0, 0.03)

    add_result_tile(slide, 1.0, 27.35, 9.8, 5.2, "DeiT-Small", "76.49 > 76.08", "Transformer result.\nEasy poster message: not CNN-specific.", "blue")
    add_result_tile(slide, 11.65, 27.35, 9.8, 5.2, "ADE20K", "42.8 > 42.4", "Segmentation result.\nBackbone features transfer better.", "teal")
    add_result_tile(slide, 22.3, 27.35, 9.8, 5.2, "Grad-CAM", "more focused attention", "More object-focused reasoning.\nLess background distraction.", "orange")

    add_textbox(slide, 1.0, 33.4, 13.2, 0.7, "SUPPORTING EVIDENCE", size=29, color="ink", bold=True, font="Aptos Display")
    add_rule(slide, 1.0, 34.2, 31.0, 0.03)

    add_card(slide, 1.0, 34.75, 10.0, 8.0, fill="paper")
    add_pill(slide, 1.35, 35.05, 2.6, 0.42, "REPRESENTATION", "sand", "ink")
    add_textbox(slide, 1.35, 35.65, 9.1, 1.0, "Better feature geometry", size=24, color="ink", bold=True, font="Aptos Display")
    add_cluster_icon(slide, 1.55, 37.0, "orange_soft", "orange", spread=False)
    add_cluster_icon(slide, 5.05, 37.0, "teal_soft", "teal", spread=True)
    add_textbox(slide, 1.45, 38.9, 3.2, 0.4, "LS: collapse", size=15, color="orange", bold=True, align="center")
    add_textbox(slide, 4.95, 38.9, 3.2, 0.4, "MaxSup: richer spread", size=15, color="teal", bold=True, align="center")
    add_textbox(slide, 1.35, 40.0, 9.0, 2.1, "Use Table 2 here.\nSame class should not collapse too much.\nDifferent classes should stay separable.", size=17, color="ink2")

    add_card(slide, 11.55, 34.75, 10.0, 8.0, fill="paper")
    add_pill(slide, 11.9, 35.05, 2.35, 0.42, "DOWNSTREAM", "blue")
    add_textbox(slide, 11.9, 35.65, 9.0, 1.0, "Not only classification", size=24, color="ink", bold=True, font="Aptos Display")
    add_small_card(slide, 11.95, 36.95, 4.25, 2.05, "Fine-grained", "CUB Birds and Stanford Cars.\nCaptures subtle details better.", "sand2")
    add_small_card(slide, 16.6, 36.95, 4.25, 2.05, "Long-tailed", "Better tradeoff across many-, medium-, and low-shot data.", "blue_soft")
    add_small_card(slide, 11.95, 39.35, 4.25, 2.05, "OOD / Corruption", "Confidence is more reliable on CIFAR-10-C.", "teal_soft")
    add_small_card(slide, 16.6, 39.35, 4.25, 2.05, "Segmentation", "Backbone features transfer better to ADE20K.", "sand2")

    add_card(slide, 22.1, 34.75, 9.9, 8.0, fill="paper")
    add_pill(slide, 22.45, 35.05, 2.15, 0.42, "APPENDIX", "orange")
    add_textbox(slide, 22.45, 35.65, 8.9, 1.0, "Use only the appendix items that help your story", size=22, color="ink", bold=True, font="Aptos Display")
    add_small_card(slide, 22.5, 37.0, 8.8, 1.35, "Table 12", "Extended transfer learning. Very poster-friendly.", "sand2")
    add_small_card(slide, 22.5, 38.6, 8.8, 1.35, "Table 13", "More baselines. Good for tough comparison questions.", "blue_soft")
    add_small_card(slide, 22.5, 40.2, 8.8, 1.35, "Table 16", "Compute efficiency. Almost no extra cost.", "teal_soft")

    footer = add_card(slide, 1.0, 43.25, 31.1, 2.0, fill="sand2", shadow=False)
    add_textbox(slide, 1.35, 43.55, 12.5, 0.6, "simple  ·  low overhead  ·  stronger representations", size=21, color="ink", bold=True, font="Aptos Display")
    add_textbox(slide, 14.5, 43.52, 8.4, 0.75, "Demo QR: put your Streamlit simulator here", size=17, color="ink2", bold=True)
    add_textbox(slide, 23.6, 43.55, 8.0, 0.8, "People stop for visuals.\nThey stay for one clear idea.", size=17, color="ink2", align="right")

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    print(build())
