"""
MaxSup A0 Poster – v4
Light minimalist modern style. 12-column grid. Strict type scale.
Output: poster/MaxSup_A0_HiFi_Poster_v4.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ────────────────────────────────────────────────────────────────────
BASE = Path("/Users/jayden/Desktop/7404 comp/project")
OUT_PATH = BASE / "poster" / "MaxSup_A0_HiFi_Poster_v4.pptx"
LOGO_PATH = BASE / "assets" / "assets_hku_logo_horizontal.png"

# ── Canvas ───────────────────────────────────────────────────────────────────
W = 33.11   # A0 portrait width  (inches)
H = 46.81   # A0 portrait height (inches)

# ── Grid (12 columns) ────────────────────────────────────────────────────────
MARGIN   = 1.1           # left / right outer margin
GUTTER   = 0.30          # gap between columns
N_COLS   = 12
COL_W    = (W - 2 * MARGIN - (N_COLS - 1) * GUTTER) / N_COLS  # ≈ 2.36"

def col_x(start_col, span):
    """Left x and width for columns start_col..(start_col+span-1), 0-indexed."""
    x = MARGIN + start_col * (COL_W + GUTTER)
    w = span * COL_W + (span - 1) * GUTTER
    return x, w

# ── Vertical rhythm ───────────────────────────────────────────────────────────
SEC_GAP  = 0.90   # space between sections
CARD_PAD = 0.45   # internal card padding

# ── Typography scale ──────────────────────────────────────────────────────────
T_TITLE   = 90    # poster title
T_SUB     = 36    # subtitle / section label
T_H1      = 52    # section header
T_H2      = 36    # card header
T_METRIC  = 88    # big result number
T_BODY    = 26    # body text
T_SMALL   = 21    # secondary body
T_CAPTION = 17    # captions / footnotes

F_DISPLAY = "Aptos Display"
F_BODY    = "Aptos"
F_MONO    = "Courier New"   # for formula lines

# ── Colour palette ────────────────────────────────────────────────────────────
PALETTE = {
    # grounds
    "bg":          "F7F6F2",   # warm off-white
    "paper":       "FFFFFF",
    "card":        "F0EEE9",   # very light warm grey for alt cards
    # ink
    "ink":         "0E1116",
    "muted":       "5B6470",
    "hairline":    "E0DDD7",
    # accent (MaxSup / right side)
    "accent":      "C1121F",   # strong crimson
    "accent_soft": "FCECEA",
    # support (LS / left side)
    "navy":        "1E3A8A",
    "navy_soft":   "E8EDF8",
    # neutral bar background
    "bar_bg":      "E8E6DF",
    # positive result
    "green":       "1B7A4E",
    "green_soft":  "E2F4EC",
    # white
    "white":       "FFFFFF",
}


def rgb(name_or_hex: str) -> RGBColor:
    v = PALETTE.get(name_or_hex, name_or_hex).replace("#", "")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


# ── Primitive helpers (kept from v3, lightly updated) ─────────────────────────

def add_shape(slide, kind, x, y, w, h, fill, line=None, line_pt=1.0, transparency=0.0):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.fill.transparency = transparency
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def add_rect(slide, x, y, w, h, fill, line=None):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line)


def add_rounded(slide, x, y, w, h, fill, line=None, radius=0.10):
    s = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    s.adjustments[0] = radius
    return s


def add_rule(slide, x, y, w, fill="hairline"):
    return add_rect(slide, x, y, w, 0.025, fill)


def add_card(slide, x, y, w, h, fill="paper", line="hairline"):
    """Flat card — no shadow (minimalist style)."""
    return add_rounded(slide, x, y, w, h, fill, line)


def add_textbox(slide, x, y, w, h, text, size=T_BODY, color="ink",
                bold=False, font=F_BODY, align="left", valign="top", margin=0.08,
                line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin * 0.5)
    tf.margin_bottom = Inches(margin * 0.5)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM
    }[valign]
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        # set line spacing via pPr
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    return box


def add_pill(slide, x, y, w, h, text, fill, text_color="white", size=15):
    s = add_rounded(slide, x, y, w, h, fill, radius=0.50)
    add_textbox(slide, x, y + 0.01, w, h - 0.02, text, size=size,
                color=text_color, bold=True, align="center", valign="middle")
    return s


# ── New v4 helpers ────────────────────────────────────────────────────────────

def hero_bar_group(slide, x, y, w, h, title, title_color,
                   values, labels, highlight_idx, arrow_color, suppress_label):
    """
    Draws a labelled bar chart for the hero section.
    highlight_idx bar gets arrow + suppress_label below.
    """
    bar_area_h = h - 1.2
    max_val    = max(values) * 1.1
    n          = len(values)
    bar_w      = (w - 0.4) / n * 0.55
    spacing    = (w - 0.4) / n
    base_y     = y + h - 0.55

    # title
    add_textbox(slide, x, y, w, 0.65, title, size=T_H2, color=title_color,
                bold=True, font=F_DISPLAY, align="center")

    for i, (lab, val) in enumerate(zip(labels, values)):
        bx     = x + 0.2 + i * spacing + (spacing - bar_w) / 2
        bar_h  = bar_area_h * (val / max_val)
        bar_y  = base_y - bar_h
        fill   = title_color if i == highlight_idx else "bar_bg"
        add_rounded(slide, bx, bar_y, bar_w, bar_h, fill, radius=0.05)
        # value label above bar
        add_textbox(slide, bx - 0.1, bar_y - 0.38, bar_w + 0.2, 0.34,
                    f"{val:.1f}", size=T_SMALL, color="ink",
                    bold=(i == highlight_idx), align="center")
        # class label below
        add_textbox(slide, bx - 0.1, base_y + 0.06, bar_w + 0.2, 0.34,
                    lab, size=T_SMALL, color="muted", align="center")

    # down arrow over highlighted bar
    hi_bx = x + 0.2 + highlight_idx * spacing + (spacing - bar_w) / 2
    add_shape(slide, MSO_SHAPE.DOWN_ARROW,
              hi_bx + bar_w * 0.1, y + 0.72, bar_w * 0.8, 0.55, arrow_color)
    # suppress label
    add_textbox(slide, x, base_y + 0.46, w, 0.42,
                suppress_label, size=T_CAPTION, color=arrow_color,
                bold=True, align="center")


def metric_tile(slide, x, y, w, h, label, pill_fill, big_number, subtitle):
    """Clean result tile: pill label, huge metric, subtitle."""
    add_card(slide, x, y, w, h, fill="paper")
    px, pw = x + CARD_PAD, 2.6
    add_pill(slide, px, y + CARD_PAD, pw, 0.46, label, pill_fill, size=16)
    add_textbox(slide, x + CARD_PAD, y + 1.25, w - 2 * CARD_PAD, 1.3,
                big_number, size=T_METRIC, color="ink", bold=True,
                font=F_DISPLAY, align="left")
    add_textbox(slide, x + CARD_PAD, y + 2.55, w - 2 * CARD_PAD, h - 2.9,
                subtitle, size=T_SMALL, color="muted", align="left")


def why_chip(slide, x, y, w, h, icon_text, label, body):
    """Small 'why it matters' feature chip."""
    add_card(slide, x, y, w, h, fill="card")
    add_textbox(slide, x + CARD_PAD, y + 0.30, w - 2 * CARD_PAD, 0.65,
                icon_text + "  " + label,
                size=T_H2, color="ink", bold=True, font=F_DISPLAY)
    add_textbox(slide, x + CARD_PAD, y + 1.05, w - 2 * CARD_PAD, h - 1.3,
                body, size=T_SMALL, color="muted")


# ── Main build ────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Background ──────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, H, "bg")

    # ── 1. HEADER  (y = 0.7 → 5.0) ──────────────────────────────────────────
    # Solid dark header band for drama
    add_rect(slide, 0, 0, W, 5.2, "ink")

    # Title line: "MaxSup" with accent dot
    cx, cw = col_x(0, 9)
    add_textbox(slide, cx, 0.55, cw, 1.8,
                "MaxSup", size=T_TITLE, color="white",
                bold=True, font=F_DISPLAY, align="left")
    # Accent dot after title
    add_shape(slide, MSO_SHAPE.OVAL,
              cx + 9.55, 0.85, 0.38, 0.38, "accent")

    add_textbox(slide, cx, 2.35, cw + 1.5, 1.1,
                "Overcoming Representation Collapse in Label Smoothing",
                size=T_SUB, color="hairline",
                bold=False, font=F_DISPLAY, align="left")

    # Author / group info (right side of header)
    rx, rw = col_x(9, 3)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH),
                                 Inches(rx), Inches(0.7), width=Inches(rw))
    add_textbox(slide, rx, 2.2, rw, 0.6,
                "Group {{GROUP_ID}}",
                size=T_SMALL, color="hairline", bold=True, align="right")
    add_textbox(slide, rx, 2.82, rw, 1.1,
                "{{NAME_1}}  ·  {{NAME_2}}  ·  {{NAME_3}}",
                size=T_CAPTION, color="muted", align="right")
    add_textbox(slide, rx, 3.7, rw, 0.55,
                "COMP7404  ·  HKU",
                size=T_CAPTION, color="muted", align="right")

    # ── 2. HOOK  (y = 5.5 → 7.0) ────────────────────────────────────────────
    hx, hw = col_x(0, 12)
    add_textbox(slide, hx, 5.55, hw, 1.15,
                "\u201cThe class that should be regularized is not always the labeled class \u2014 "
                "it\u2019s the class the model currently trusts most.\u201d",
                size=T_H2, color="ink", bold=False,
                font=F_DISPLAY, align="center")
    add_rule(slide, hx, 6.82, hw)

    # ── 3. HERO  (y = 7.3 → 20.5) ───────────────────────────────────────────
    hero_y    = 7.3
    hero_h    = 13.0
    hero_bot  = hero_y + hero_h

    # Section label
    add_textbox(slide, hx, hero_y, 6.0, 0.6,
                "THE CORE INSIGHT", size=T_SMALL,
                color="muted", bold=True, font=F_BODY, align="left")

    # ── Left panel: "What goes wrong with LS" (cols 0-4) ────────────────────
    lx, lw = col_x(0, 5)
    lcard_y = hero_y + 0.75
    lcard_h = hero_h - 0.75
    add_card(slide, lx, lcard_y, lw, lcard_h, fill="navy_soft", line="hairline")

    add_pill(slide, lx + CARD_PAD, lcard_y + CARD_PAD,
             3.2, 0.46, "MISCLASSIFIED SAMPLE", "navy", size=14)

    # Stylised cat icon (procedural)
    cx_cat = lx + lw / 2
    cy_cat = lcard_y + 2.1
    # body
    add_shape(slide, MSO_SHAPE.OVAL, cx_cat - 1.1, cy_cat - 0.85, 2.2, 1.7, "bar_bg")
    # head
    add_shape(slide, MSO_SHAPE.OVAL, cx_cat - 0.7, cy_cat - 1.85, 1.4, 1.4, "bar_bg")
    # ears (triangles via right-triangle shape)
    add_shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, cx_cat - 0.65, cy_cat - 2.45, 0.35, 0.6, "bar_bg")
    add_shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, cx_cat + 0.30, cy_cat - 2.45, 0.35, 0.6, "bar_bg")
    # label below icon
    add_textbox(slide, lx + CARD_PAD, lcard_y + 1.45, lw - 2*CARD_PAD, 0.5,
                "GT = cat  (but model predicts fox)",
                size=T_CAPTION, color="muted", align="center")

    # Logit table
    table_y = lcard_y + 4.0
    row_h   = 0.82
    rows    = [("fox", 2.8, True), ("cat", 1.0, False), ("dog", 0.4, False)]
    for i, (cls, val, is_wrong) in enumerate(rows):
        ry = table_y + i * row_h
        bg = "accent_soft" if is_wrong else "paper"
        add_rounded(slide, lx + CARD_PAD, ry, lw - 2*CARD_PAD, row_h - 0.08, bg, "hairline")
        # class name
        add_textbox(slide, lx + CARD_PAD + 0.15, ry + 0.12,
                    1.5, row_h - 0.24,
                    cls, size=T_BODY, color="accent" if is_wrong else "ink",
                    bold=is_wrong, align="left", valign="middle")
        # bar
        bar_max_w = lw - 2*CARD_PAD - 2.3
        bar_ww    = max(0.15, bar_max_w * val / 3.2)
        add_rounded(slide, lx + CARD_PAD + 1.75, ry + 0.22,
                    bar_ww, row_h - 0.44,
                    "accent" if is_wrong else "bar_bg")
        # value
        add_textbox(slide, lx + lw - CARD_PAD - 0.85, ry + 0.12,
                    0.7, row_h - 0.24,
                    f"{val:.1f}", size=T_BODY, color="accent" if is_wrong else "muted",
                    bold=is_wrong, align="right", valign="middle")

    add_textbox(slide, lx + CARD_PAD, table_y + 3 * row_h + 0.12,
                lw - 2*CARD_PAD, 0.55,
                "Model picks fox  (wrong)",
                size=T_CAPTION, color="accent", bold=True, align="center")

    # ── Right panel: LS vs MaxSup bar charts (cols 5-11) ────────────────────
    rx2, rw2 = col_x(5, 7)
    rcard_y  = lcard_y
    rcard_h  = lcard_h
    add_card(slide, rx2, rcard_y, rw2, rcard_h, fill="paper", line="hairline")

    # "The difference" headline
    add_textbox(slide, rx2 + CARD_PAD, rcard_y + CARD_PAD, rw2 - 2*CARD_PAD, 0.7,
                "How each method responds",
                size=T_H2, color="ink", bold=True, font=F_DISPLAY)

    half_w  = (rw2 - 2*CARD_PAD - GUTTER) / 2
    lsbar_x = rx2 + CARD_PAD
    msbar_x = rx2 + CARD_PAD + half_w + GUTTER

    hero_bar_group(slide, lsbar_x, rcard_y + 1.55, half_w, 8.5,
                   "Label Smoothing", "navy",
                   [1.0, 2.8, 0.4], ["cat", "fox", "dog"],
                   highlight_idx=0,      # LS pushes down cat (GT)
                   arrow_color="navy",
                   suppress_label="\u2193 suppresses GT cat  \u2717")

    hero_bar_group(slide, msbar_x, rcard_y + 1.55, half_w, 8.5,
                   "MaxSup", "accent",
                   [1.0, 2.8, 0.4], ["cat", "fox", "dog"],
                   highlight_idx=1,      # MaxSup pushes down fox (wrong top-1)
                   arrow_color="accent",
                   suppress_label="\u2193 suppresses wrong fox  \u2713")

    # Formula comparison strip
    fstrip_y = rcard_y + 10.35
    add_rounded(slide, rx2 + CARD_PAD, fstrip_y,
                rw2 - 2*CARD_PAD, 1.55, "card", "hairline")
    add_textbox(slide, rx2 + CARD_PAD + 0.2, fstrip_y + 0.12,
                rw2 - 2*CARD_PAD - 0.4, 0.5,
                "L_LS = \u03b1 ( z_gt \u2212 mean(z) )",
                size=T_SMALL, color="navy", bold=True, font=F_MONO, align="left")
    add_textbox(slide, rx2 + CARD_PAD + 0.2, fstrip_y + 0.68,
                rw2 - 2*CARD_PAD - 0.4, 0.5,
                "L_MaxSup = \u03b1 ( z_max \u2212 mean(z) )  \u2190 one symbol changed",
                size=T_SMALL, color="accent", bold=True, font=F_MONO, align="left")

    # Hero caption
    add_textbox(slide, hx, hero_bot + 0.12, hw, 0.7,
                "LS penalizes the labeled class even when the model is already wrong. "
                "MaxSup penalizes whichever class the model currently trusts most \u2014 "
                "a consistent signal in both correct and wrong predictions.",
                size=T_SMALL, color="muted", align="center")

    add_rule(slide, hx, hero_bot + 0.95, hw)

    # ── 4. PROBLEM | METHOD  (y = 22.0 → 30.0) ──────────────────────────────
    sec4_y = hero_bot + 1.15

    # Section header
    add_textbox(slide, hx, sec4_y, 8.0, 0.75,
                "PROBLEM & METHOD",
                size=T_H1, color="ink", bold=True, font=F_DISPLAY)

    col4_y = sec4_y + 1.0

    # LEFT card: Problem / LS decomposition  (cols 0-5)
    plx, plw = col_x(0, 6)
    plh = 7.5
    add_card(slide, plx, col4_y, plw, plh, fill="navy_soft", line="hairline")

    add_pill(slide, plx + CARD_PAD, col4_y + CARD_PAD,
             2.8, 0.46, "THE PROBLEM", "navy", size=14)
    add_textbox(slide, plx + CARD_PAD, col4_y + 1.2, plw - 2*CARD_PAD, 1.05,
                "LS decomposes into two terms:",
                size=T_BODY, color="ink", bold=True, font=F_DISPLAY)

    # Decomposition boxes
    for i, (label, body, c) in enumerate([
        ("Regularization term",
         "Acts on logits below z_gt.\nReduces overconfidence. Good.",
         "green"),
        ("Error amplification term",
         "Acts on logits above z_gt.\nWhen model is wrong, widens the gap. Bad.",
         "accent"),
    ]):
        iy = col4_y + 2.45 + i * 2.3
        bg = "green_soft" if c == "green" else "accent_soft"
        add_rounded(slide, plx + CARD_PAD, iy, plw - 2*CARD_PAD, 1.85, bg, "hairline")
        add_textbox(slide, plx + CARD_PAD + 0.15, iy + 0.18,
                    plw - 2*CARD_PAD - 0.3, 0.5,
                    label, size=T_BODY, color=c, bold=True, font=F_DISPLAY)
        add_textbox(slide, plx + CARD_PAD + 0.15, iy + 0.72,
                    plw - 2*CARD_PAD - 0.3, 0.9,
                    body, size=T_SMALL, color="ink")

    add_textbox(slide, plx + CARD_PAD, col4_y + plh - 0.65,
                plw - 2*CARD_PAD, 0.48,
                "On misclassified samples, the bad term dominates.",
                size=T_SMALL, color="muted", bold=True)

    # RIGHT card: Method  (cols 6-11)
    pmx, pmw = col_x(6, 6)
    pmh = plh
    add_card(slide, pmx, col4_y, pmw, pmh, fill="accent_soft", line="hairline")

    add_pill(slide, pmx + CARD_PAD, col4_y + CARD_PAD,
             2.5, 0.46, "THE FIX", "accent", size=14)
    add_textbox(slide, pmx + CARD_PAD, col4_y + 1.2, pmw - 2*CARD_PAD, 0.7,
                "One substitution eliminates the bad term:",
                size=T_BODY, color="ink", bold=True, font=F_DISPLAY)

    # Big formula highlight
    fbox_y = col4_y + 2.1
    add_rounded(slide, pmx + CARD_PAD, fbox_y, pmw - 2*CARD_PAD, 1.2, "paper", "hairline")
    add_textbox(slide, pmx + CARD_PAD + 0.2, fbox_y + 0.1,
                pmw - 2*CARD_PAD - 0.4, 0.46,
                "z_gt   \u279c   z_max",
                size=T_H1, color="accent", bold=True, font=F_DISPLAY, align="center")
    add_textbox(slide, pmx + CARD_PAD + 0.2, fbox_y + 0.65,
                pmw - 2*CARD_PAD - 0.4, 0.38,
                "replace ground-truth logit with top-1 logit",
                size=T_CAPTION, color="muted", align="center")

    # Two behaviour cards
    for i, (title, body, c, bg) in enumerate([
        ("Correct prediction",
         "z_max = z_gt\nActs exactly like confidence regularization.",
         "green", "green_soft"),
        ("Wrong prediction",
         "z_max \u2260 z_gt\nDirectly suppresses the wrong dominant class.",
         "accent", "accent_soft"),
    ]):
        iy = col4_y + 3.6 + i * 2.05
        add_rounded(slide, pmx + CARD_PAD, iy, pmw - 2*CARD_PAD, 1.65, bg, "hairline")
        add_textbox(slide, pmx + CARD_PAD + 0.15, iy + 0.18,
                    pmw - 2*CARD_PAD - 0.3, 0.48,
                    title, size=T_BODY, color=c, bold=True, font=F_DISPLAY)
        add_textbox(slide, pmx + CARD_PAD + 0.15, iy + 0.7,
                    pmw - 2*CARD_PAD - 0.3, 0.75,
                    body, size=T_SMALL, color="ink")

    add_textbox(slide, pmx + CARD_PAD, col4_y + pmh - 0.65,
                pmw - 2*CARD_PAD, 0.48,
                "Consistent regularization signal regardless of correctness.",
                size=T_SMALL, color="muted", bold=True)

    add_rule(slide, hx, col4_y + plh + 0.6, hw)

    # ── 5. RESULTS  (y ≈ 32.0 → 38.5) ───────────────────────────────────────
    sec5_y = col4_y + plh + 0.9

    add_textbox(slide, hx, sec5_y, 8.0, 0.75,
                "RESULTS",
                size=T_H1, color="ink", bold=True, font=F_DISPLAY)

    tile_y = sec5_y + 1.0
    tile_h = 6.0
    tw, _ = col_x(0, 4)   # each tile = 4 cols
    tw    = (hw - 2 * GUTTER) / 3

    for i, (label, pill_c, big, sub) in enumerate([
        ("DeiT-Small  ImageNet",  "navy",
         "76.49%",
         "+0.41 over LS\n+2.10 over Baseline\nWorks on Transformers, not just CNNs."),
        ("ADE20K  Segmentation",  "green",
         "42.8 mIoU",
         "+0.4 over LS\nBackbone features transfer better\nto dense prediction tasks."),
        ("Grad-CAM  Attention",   "accent",
         "More focused",
         "LS attends to background cues.\nMaxSup attends to the object body.\nBetter visual interpretability."),
    ]):
        tx = hx + i * (tw + GUTTER)
        metric_tile(slide, tx, tile_y, tw, tile_h, label, pill_c, big, sub)

    add_rule(slide, hx, tile_y + tile_h + 0.6, hw)

    # ── 6. WHY IT MATTERS  (y ≈ 40.0 → 43.5) ───────────────────────────────
    sec6_y = tile_y + tile_h + 0.9

    add_textbox(slide, hx, sec6_y, 10.0, 0.75,
                "WHY IT MATTERS",
                size=T_H1, color="ink", bold=True, font=F_DISPLAY)

    chip_y = sec6_y + 1.0
    chip_h = 2.85
    cw     = (hw - 3 * GUTTER) / 4

    chips = [
        ("[1]", "Plug-in",
         "No architecture change. Drop-in loss replacement."),
        ("[2]", "Consistent",
         "Same correct signal on all samples, right or wrong."),
        ("[3]", "Transferable",
         "Better features for segmentation, fine-grained, and OOD tasks."),
        ("[4]", "Interpretable",
         "Model attends to the right regions (Grad-CAM evidence)."),
    ]
    for i, (icon, label, body) in enumerate(chips):
        cx2 = hx + i * (cw + GUTTER)
        why_chip(slide, cx2, chip_y, cw, chip_h, icon, label, body)

    add_rule(slide, hx, chip_y + chip_h + 0.45, hw)

    # ── 7. FOOTER  (y ≈ 44.0 → 46.5) ────────────────────────────────────────
    footer_y = chip_y + chip_h + 0.65

    add_rect(slide, 0, footer_y, W, H - footer_y, "ink")

    # Paper reference
    fx, fw = col_x(0, 7)
    add_textbox(slide, fx, footer_y + 0.35, fw, 0.52,
                "MaxSup: Overcoming Representation Collapse in Label Smoothing  |  arXiv 2502.15798",
                size=T_CAPTION, color="hairline", align="left")
    add_textbox(slide, fx, footer_y + 0.92, fw, 0.52,
                "Limitation: low-shot class performance remains an open challenge.",
                size=T_CAPTION, color="muted", align="left")

    # QR placeholder
    qx, qw = col_x(8, 2)
    add_rounded(slide, qx, footer_y + 0.25, 1.8, 1.8, "bar_bg", line="hairline")
    add_textbox(slide, qx, footer_y + 0.25, 1.8, 1.8,
                "Demo\nQR", size=T_SMALL, color="muted",
                align="center", valign="middle")

    # Tagline
    tx2, tw2 = col_x(10, 2)
    add_textbox(slide, tx2, footer_y + 0.35, tw2, 1.4,
                "Simple.\nFocused.\nBetter.",
                size=T_BODY, color="white", bold=True,
                font=F_DISPLAY, align="right")

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    return OUT_PATH


if __name__ == "__main__":
    build()
